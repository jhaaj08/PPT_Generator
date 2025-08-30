from flask import Flask, request, jsonify, send_file, render_template_string, send_from_directory
from flask_cors import CORS
import os
import tempfile
import json
import time
from werkzeug.utils import secure_filename
from improved_ppt_processor import ImprovedPPTProcessor
from llm_integration import LLMIntegration
from pptx import Presentation
from error_handler import (
    safe_api_call, 
    validate_file_upload, 
    validate_text_input,
    validate_api_key,
    PPTGeneratorError,
    FileValidationError,
    LLMError,
    TemplateError,
    ProgressTracker,
    create_error_response,
    logger
)
from retry_handler import retry_file_operation, RetryableOperation, RetryConfigs
import traceback

app = Flask(__name__, static_folder='static')
CORS(app)

# Configuration
UPLOAD_FOLDER = 'uploads'
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
ALLOWED_EXTENSIONS = {'pptx', 'potx'}

# Ensure upload directory exists
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

@app.route('/')
def index():
    return send_from_directory(app.static_folder, 'index.html')

@app.route('/test-manifest')
def test_manifest():
    return send_from_directory(app.static_folder, 'test-manifest.html')

@app.route('/test-template-extraction')
def test_template_extraction():
    return send_from_directory(app.static_folder, 'test-template-extraction.html')

@app.route('/test-single-slide')
def test_single_slide():
    return send_from_directory(app.static_folder, 'test-single-slide.html')

@app.route('/test-full-flow')
def test_full_flow():
    return send_from_directory(app.static_folder, 'test-full-flow.html')

@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_from_directory(UPLOAD_FOLDER, filename)

@app.route('/<path:filename>')
def static_files(filename):
    return send_from_directory(app.static_folder, filename)

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'healthy', 'message': 'PPT Generator API is running'})

@app.route('/api/generate-presentation', methods=['POST'])
def generate_presentation():
    try:
        # Validate request
        if 'template' not in request.files:
            return jsonify({'error': 'No template file provided'}), 400
        
        template_file = request.files['template']
        if template_file.filename == '':
            return jsonify({'error': 'No template file selected'}), 400
        
        if not allowed_file(template_file.filename):
            return jsonify({'error': 'Invalid file type. Only .pptx and .potx files are allowed'}), 400
        
        # Get form data
        input_text = request.form.get('text', '').strip()
        guidance = request.form.get('guidance', '').strip()
        api_key = request.form.get('api_key', '').strip()
        llm_provider = request.form.get('llm_provider', 'openai').strip()
        include_speaker_notes = request.form.get('include_speaker_notes', 'false').lower() == 'true'
        
        if not input_text:
            return jsonify({'error': 'Input text is required'}), 400
        
        if not api_key:
            return jsonify({'error': 'API key is required'}), 400
        
        # Save uploaded template
        template_filename = secure_filename(template_file.filename)
        template_path = os.path.join(UPLOAD_FOLDER, template_filename)
        template_file.save(template_path)
        
        try:
            print(f"🚀 Starting manifest-based presentation generation...")
            print(f"   Speaker notes enabled: {include_speaker_notes}")
            
            # Initialize processors
            ppt_processor = ImprovedPPTProcessor()
            llm_integration = LLMIntegration(api_key, llm_provider)
            
            # Step 1: Extract raw template data
            raw_template_data = ppt_processor.extract_raw_template_data(template_path)
            print(f"🔍 Raw template data extracted")
            
            # Step 2: Generate manifest using LLM
            manifest = llm_integration.generate_template_manifest(raw_template_data)
            print(f"🧠 Template manifest generated with {len(manifest.get('layouts', []))} layout rules")
            
            # Step 3: Structure the content using LLM
            slide_structure = llm_integration.structure_text_to_slides(input_text, guidance)
            print(f"📋 Content structured into {len(slide_structure.get('slides', []))} slides")
            print(f"🔍 Debug: First slide structure = {slide_structure.get('slides', [{}])[0] if slide_structure.get('slides') else 'No slides found'}")
            
            # Step 4: Generate presentation using manifest
            prs = ppt_processor.generate_presentation_with_manifest(slide_structure, template_path, manifest)
            
            # Step 5: Save the presentation
            timestamp = int(time.time())
            output_filename = f'generated_presentation_{timestamp}.pptx'
            output_path = os.path.join(UPLOAD_FOLDER, output_filename)
            prs.save(output_path)
            
            print(f"✅ Manifest-based presentation generated successfully")
            
            # Return the generated file
            return send_file(
                output_path,
                as_attachment=True,
                download_name=output_filename,
                mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            
        finally:
            # Cleanup uploaded file
            if os.path.exists(template_path):
                os.remove(template_path)
        
    except Exception as e:
        print(f"Error generating presentation: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Failed to generate presentation: {str(e)}'}), 500

@app.route('/api/generate-manifest', methods=['POST'])
def generate_manifest():
    """Generate template manifest using LLM analysis"""
    try:
        # Validate request
        if 'template' not in request.files:
            return jsonify({'error': 'No template file provided'}), 400
        
        template_file = request.files['template']
        if template_file.filename == '':
            return jsonify({'error': 'No template file selected'}), 400
        
        if not allowed_file(template_file.filename):
            return jsonify({'error': 'Invalid file type. Only .pptx and .potx files are allowed'}), 400
        
        # Get form data
        api_key = request.form.get('api_key', '').strip()
        llm_provider = request.form.get('llm_provider', 'openai').strip()
        
        if not api_key:
            return jsonify({'error': 'API key is required'}), 400
        
        # Save uploaded template
        template_filename = secure_filename(template_file.filename)
        template_path = os.path.join(UPLOAD_FOLDER, template_filename)
        template_file.save(template_path)
        
        try:
            print(f"🔍 Generating manifest for template: {template_filename}")
            
            # Initialize processors
            ppt_processor = ImprovedPPTProcessor()
            llm_integration = LLMIntegration(api_key, llm_provider)
            
            # Step 1: Extract raw template data (deterministic)
            raw_data = ppt_processor.extract_raw_template_data(template_path)
            
            # Step 2: Use LLM to generate structured manifest
            manifest = llm_integration.generate_template_manifest(raw_data)
            
            print(f"✅ Manifest generated successfully")
            
            # Return both raw data and manifest for comparison
            return jsonify({
                'raw_template_data': raw_data,
                'llm_manifest': manifest,
                'template_filename': template_filename
            })
            
        finally:
            # Cleanup uploaded file
            if os.path.exists(template_path):
                os.remove(template_path)
        
    except Exception as e:
        print(f"Error generating manifest: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Failed to generate manifest: {str(e)}'}), 500

@app.route('/api/test-template-extraction', methods=['POST'])
def test_template_extraction_api():
    """Test template extraction without LLM - core extraction only"""
    try:
        # Validate request
        if 'template' not in request.files:
            return jsonify({'error': 'No template file provided'}), 400
        
        template_file = request.files['template']
        if template_file.filename == '':
            return jsonify({'error': 'No template file selected'}), 400
        
        if not allowed_file(template_file.filename):
            return jsonify({'error': 'Invalid file type. Only .pptx and .potx files are allowed'}), 400
        
        # Save uploaded template
        template_filename = secure_filename(template_file.filename)
        template_path = os.path.join(UPLOAD_FOLDER, template_filename)
        template_file.save(template_path)
        
        try:
            print(f"🔍 Testing template extraction: {template_filename}")
            
            # Initialize processor
            ppt_processor = ImprovedPPTProcessor()
            
            # Extract template data (this is what we're testing)
            raw_data = ppt_processor.extract_raw_template_data(template_path)
            
            print(f"✅ Template extraction successful")
            
            return jsonify({
                'success': True,
                'template_filename': template_filename,
                'basic_info': {
                    'slide_width': raw_data.get('slide_size', {}).get('width_emu', 0),
                    'slide_height': raw_data.get('slide_size', {}).get('height_emu', 0),
                    'total_layouts': len(raw_data.get('layouts', [])),
                    'total_images': len(raw_data.get('images', []))
                },
                'layouts': raw_data.get('layouts', []),
                'theme': raw_data.get('theme', {}),
                'images': raw_data.get('images', []),
                'slide_size': raw_data.get('slide_size', {})
            })
            
        finally:
            # Cleanup uploaded file
            if os.path.exists(template_path):
                os.remove(template_path)
        
    except Exception as e:
        print(f"Error testing template extraction: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Failed to extract template: {str(e)}'}), 500

@app.route('/api/test-single-slide', methods=['POST'])
def test_single_slide_api():
    """Test single slide generation - LLM content + template layout matching"""
    try:
        # Validate request
        if 'template' not in request.files:
            return jsonify({'error': 'No template file provided'}), 400
        
        template_file = request.files['template']
        if template_file.filename == '':
            return jsonify({'error': 'No template file selected'}), 400
        
        if not allowed_file(template_file.filename):
            return jsonify({'error': 'Invalid file type. Only .pptx and .potx files are allowed'}), 400
        
        # Get form data
        input_text = request.form.get('text', '').strip()
        api_key = request.form.get('api_key', '').strip()
        llm_provider = request.form.get('llm_provider', 'openai').strip()
        
        if not input_text:
            return jsonify({'error': 'Input text is required'}), 400
        
        if not api_key:
            return jsonify({'error': 'API key is required'}), 400
        
        # Save uploaded template
        template_filename = secure_filename(template_file.filename)
        template_path = os.path.join(UPLOAD_FOLDER, template_filename)
        template_file.save(template_path)
        
        try:
            print(f"🧪 Testing single slide generation: {template_filename}")
            
            # Step 1: Initialize processors
            ppt_processor = ImprovedPPTProcessor()
            llm_integration = LLMIntegration(api_key, llm_provider)
            
            # Step 2: Generate slide content with LLM
            slide_structure = llm_integration.structure_text_to_slides(input_text, "Create a professional slide")
            
            if not slide_structure.get('slides'):
                return jsonify({'error': 'LLM failed to generate slide content'}), 500
            
            # Take only the first slide for testing
            first_slide = slide_structure['slides'][0]
            
            # Step 3: Extract template data and generate manifest
            raw_template_data = ppt_processor.extract_raw_template_data(template_path)
            manifest = llm_integration.generate_template_manifest(raw_template_data)
            
            # Step 4: Generate single slide
            prs = Presentation(template_path)
            
            # Clear existing slides
            slide_count = len(prs.slides)
            for i in range(slide_count - 1, -1, -1):
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
            
            # Create the test slide
            slide_idx = 0
            layout_choice = ppt_processor._resolve_layout_from_manifest(first_slide, manifest, slide_idx, prs)
            layout_matched = layout_choice is not None
            
            if not layout_choice:
                layout_choice = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            
            slide = prs.slides.add_slide(layout_choice)
            
            # Apply content and styling
            ppt_processor._apply_manifest_styling(slide, first_slide, manifest)
            ppt_processor._add_manifest_assets(slide, manifest, slide_idx)
            
            # Step 5: Save and verify result
            timestamp = int(time.time())
            output_filename = f'single_slide_test_{timestamp}.pptx'
            output_path = os.path.join(UPLOAD_FOLDER, output_filename)
            prs.save(output_path)
            
            # Verify the result
            verify_prs = Presentation(output_path)
            verify_slide = verify_prs.slides[0] if len(verify_prs.slides) > 0 else None
            
            title_applied = False
            content_applied = False
            final_title = ""
            final_content = ""
            content_shapes = 0
            
            if verify_slide:
                # Check title
                if hasattr(verify_slide.shapes, 'title') and verify_slide.shapes.title:
                    final_title = verify_slide.shapes.title.text
                    title_applied = bool(final_title.strip())
                
                # Check content
                for shape in verify_slide.shapes:
                    if (hasattr(shape, 'text_frame') and 
                        shape.text_frame and 
                        shape != getattr(verify_slide.shapes, 'title', None)):
                        if shape.text_frame.text.strip():
                            content_shapes += 1
                            final_content = shape.text_frame.text
                            content_applied = True
                            break
            
            print(f"✅ Single slide test completed")
            
            return jsonify({
                'success': title_applied and content_applied,
                'llm_success': True,
                'layout_matched': layout_matched,
                'content_applied': content_applied,
                'title_applied': title_applied,
                'generated_slide': first_slide,
                'layout_used': layout_choice.name,
                'available_layouts': len(prs.slide_layouts),
                'final_title': final_title,
                'final_content': final_content[:100] + "..." if len(final_content) > 100 else final_content,
                'content_shapes': content_shapes,
                'template_filename': template_filename,
                'download_url': f'/uploads/{output_filename}' if os.path.exists(output_path) else None,
                'file_size': os.path.getsize(output_path) if os.path.exists(output_path) else 0
            })
            
        finally:
            # Cleanup uploaded template (but keep generated file for download)
            if os.path.exists(template_path):
                os.remove(template_path)
        
    except Exception as e:
        print(f"Error testing single slide generation: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Failed to generate single slide: {str(e)}'}), 500

@app.route('/api/generate-presentation-with-preview', methods=['POST'])
@safe_api_call
def generate_presentation_with_preview():
    """Generate full presentation with preview data"""
    # Initialize progress tracker
    progress = ProgressTracker()
    progress.add_step("validate_input", "Validating input data")
    progress.add_step("extract_template", "Analyzing template")  
    progress.add_step("generate_manifest", "Creating style manifest")
    progress.add_step("structure_content", "Structuring content")
    progress.add_step("generate_presentation", "Generating presentation")
    progress.add_step("prepare_preview", "Preparing preview data")
    
    progress.start_step(0)
    
    try:
        # Get files and form data first
        template_file = request.files.get('template')
        input_text = request.form.get('text', '').strip()
        guidance = request.form.get('guidance', '').strip()
        api_key = request.form.get('api_key', '').strip()
        llm_provider = request.form.get('llm_provider', 'openai').strip()
        include_speaker_notes = request.form.get('include_speaker_notes', 'false').lower() == 'true'
        
        # Validate all inputs using error handler
        validate_file_upload(template_file)
        validate_text_input(input_text)
        validate_api_key(api_key, llm_provider)
        
        progress.complete_step(0)
        
        # Save uploaded template
        template_filename = secure_filename(template_file.filename)
        template_path = os.path.join(UPLOAD_FOLDER, template_filename)
        template_file.save(template_path)
        
        try:
            print(f"🚀 Starting full presentation generation with preview...")
            print(f"   Speaker notes enabled: {include_speaker_notes}")
            
            # Step 1: Initialize processors
            ppt_processor = ImprovedPPTProcessor()
            llm_integration = LLMIntegration(api_key, llm_provider)
            
            # Step 2: Extract raw template data
            raw_template_data = ppt_processor.extract_raw_template_data(template_path)
            print(f"🔍 Raw template data extracted")
            
            # Step 3: Generate manifest using LLM
            manifest = llm_integration.generate_template_manifest(raw_template_data)
            print(f"🧠 Template manifest generated with {len(manifest.get('layouts', []))} layout rules")
            
            # Step 4: Structure the content using LLM
            slide_structure = llm_integration.structure_text_to_slides(input_text, guidance)
            print(f"📋 Content structured into {len(slide_structure.get('slides', []))} slides")
            
            # Step 5: Generate presentation using manifest
            prs = ppt_processor.generate_presentation_with_manifest(slide_structure, template_path, manifest)
            
            # Step 6: Save the presentation
            timestamp = int(time.time())
            output_filename = f'presentation_{timestamp}.pptx'
            output_path = os.path.join(UPLOAD_FOLDER, output_filename)
            prs.save(output_path)
            
            print(f"✅ Full presentation generated successfully")
            
            # Step 7: Prepare preview data
            slides_data = slide_structure.get('slides', [])
            layouts_used = len(manifest.get('layouts', []))
            template_images = len(raw_template_data.get('images', []))
            
            # Count slides with speaker notes
            slides_with_notes = sum(1 for slide in slides_data if slide.get('speaker_notes'))
            
            return jsonify({
                'success': True,
                'presentation_title': slide_structure.get('presentation_title', 'Generated Presentation'),
                'total_slides': len(slides_data),
                'slides': slides_data,
                'layouts_used': layouts_used,
                'template_images': template_images,
                'template_filename': template_filename,
                'download_url': f'/uploads/{output_filename}',
                'filename': output_filename,
                'file_size': os.path.getsize(output_path) if os.path.exists(output_path) else 0,
                'speaker_notes_enabled': include_speaker_notes,
                'slides_with_notes': slides_with_notes,
                'manifest_summary': {
                    'colors_extracted': len(manifest.get('theme', {}).get('palette', {})),
                    'fonts_extracted': len(manifest.get('theme', {}).get('fonts', {})),
                    'layouts_analyzed': len(manifest.get('layouts', []))
                }
            })
            
        finally:
            # Cleanup uploaded template (but keep generated file for download)
            if os.path.exists(template_path):
                os.remove(template_path)
        
    except Exception as e:
        print(f"Error generating presentation with preview: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Failed to generate presentation: {str(e)}'}), 500

@app.route('/api/supported-providers', methods=['GET'])
def get_supported_providers():
    return jsonify({
        'providers': [
            {'id': 'openai', 'name': 'OpenAI (GPT-4)'},
            {'id': 'anthropic', 'name': 'Anthropic (Claude)'},
            {'id': 'gemini', 'name': 'Google Gemini'}
        ]
    })

@app.errorhandler(413)
def too_large(e):
    return jsonify({'error': 'File too large. Maximum size is 50MB'}), 413

if __name__ == '__main__':
    import os
    
    # Create uploads directory if it doesn't exist
    os.makedirs(UPLOAD_FOLDER, exist_ok=True)
    
    # Get port from environment variable for deployment
    port = int(os.environ.get('PORT', 8080))
    debug = os.environ.get('FLASK_ENV') != 'production'
    
    print("🚀 Starting PPT Generator")
    print("📋 Make sure you have your LLM API key ready!")
    print("✨ Ready to transform your text into presentations!")
    
    app.run(host='0.0.0.0', port=port, debug=debug)

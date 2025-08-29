#!/usr/bin/env python3
"""
Improved PPT Processor that better preserves template styling
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import tempfile
import json
from PIL import Image
import io

class ImprovedPPTProcessor:
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()
    
    def analyze_template(self, template_path):
        """Enhanced template analysis that extracts more styling information"""
        try:
            prs = Presentation(template_path)
            analysis = {
                'layouts': [],
                'colors': [],
                'fonts': [],
                'images': [],
                'backgrounds': [],
                'slide_dimensions': {
                    'width': prs.slide_width,
                    'height': prs.slide_height
                },
                'master_slides': [],
                'theme_colors': []
            }
            
            print(f"🔍 Analyzing template with {len(prs.slides)} slides, {len(prs.slide_layouts)} layouts")
            
            # Analyze slide layouts with more detail
            for i, layout in enumerate(prs.slide_layouts):
                layout_info = {
                    'index': i,
                    'name': layout.name,
                    'placeholders': [],
                    'background': None
                }
                
                # Analyze placeholders
                for placeholder in layout.placeholders:
                    try:
                        placeholder_info = {
                            'index': placeholder.placeholder_format.idx,
                            'type': str(placeholder.placeholder_format.type),
                            'name': placeholder.name,
                            'left': placeholder.left,
                            'top': placeholder.top,
                            'width': placeholder.width,
                            'height': placeholder.height
                        }
                        layout_info['placeholders'].append(placeholder_info)
                    except:
                        pass
                
                # Check layout background
                try:
                    if layout.background.fill.type:
                        layout_info['background'] = str(layout.background.fill.type)
                except:
                    pass
                
                analysis['layouts'].append(layout_info)
            
            # Analyze master slides and themes
            for i, master in enumerate(prs.slide_masters):
                master_info = {
                    'name': master.name,
                    'colors': [],
                    'fonts': [],
                    'background': None
                }
                
                # Extract theme colors more thoroughly
                try:
                    theme = master.theme
                    if theme and theme.color_scheme:
                        for j in range(12):
                            try:
                                color = theme.color_scheme[j]
                                if color and hasattr(color, 'rgb'):
                                    color_hex = f"#{color.rgb:06x}"
                                    master_info['colors'].append(color_hex)
                                    if color_hex not in analysis['theme_colors']:
                                        analysis['theme_colors'].append(color_hex)
                            except:
                                pass
                except:
                    pass
                
                # Extract fonts from master
                for shape in master.shapes:
                    if hasattr(shape, 'text_frame'):
                        self._extract_text_styles(shape, analysis)
                
                analysis['master_slides'].append(master_info)
            
            # Analyze existing slides for actual styling
            for slide_idx, slide in enumerate(prs.slides):
                print(f"   Analyzing slide {slide_idx + 1}")
                self._extract_slide_styles_enhanced(slide, analysis)
            
            print(f"✅ Template analysis complete:")
            print(f"   - Colors found: {len(analysis['colors'])}")
            print(f"   - Theme colors: {len(analysis['theme_colors'])}")
            print(f"   - Fonts found: {len(analysis['fonts'])}")
            print(f"   - Images found: {len(analysis['images'])}")
            
            return analysis
            
        except Exception as e:
            print(f"❌ Error analyzing template: {str(e)}")
            import traceback
            traceback.print_exc()
            return self._get_default_analysis()
    
    def _extract_text_styles(self, shape, analysis):
        """Extract text styling information from a shape"""
        try:
            if not hasattr(shape, 'text_frame'):
                return
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Extract font information
                    if run.font.name:
                        font_info = {
                            'name': run.font.name,
                            'size': run.font.size.pt if run.font.size else 18,
                            'bold': run.font.bold or False,
                            'italic': run.font.italic or False
                        }
                        
                        # Avoid duplicates
                        if not any(f['name'] == font_info['name'] for f in analysis['fonts']):
                            analysis['fonts'].append(font_info)
                    
                    # Extract color information
                    try:
                        if run.font.color.rgb:
                            color_hex = f"#{run.font.color.rgb:06x}"
                            if color_hex not in analysis['colors']:
                                analysis['colors'].append(color_hex)
                    except:
                        pass
        except:
            pass
    
    def _extract_slide_styles_enhanced(self, slide, analysis):
        """Enhanced slide style extraction"""
        for shape in slide.shapes:
            # Extract text styles
            self._extract_text_styles(shape, analysis)
            
            # Extract images with better error handling
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_data = {
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'image_data': shape.image.blob
                    }
                    analysis['images'].append(image_data)
                    print(f"   Found image: {shape.width}x{shape.height} ({len(shape.image.blob)} bytes)")
                except Exception as e:
                    print(f"   Image extraction failed: {str(e)}")
            
            # Extract shape fills and colors
            try:
                if hasattr(shape, 'fill') and shape.fill.solid():
                    color = shape.fill.fore_color
                    if color.rgb:
                        color_hex = f"#{color.rgb:06x}"
                        if color_hex not in analysis['colors']:
                            analysis['colors'].append(color_hex)
            except:
                pass
    
    def generate_presentation_with_manifest(self, slide_structure, template_path, manifest):
        """Generate presentation using LLM-generated manifest"""
        try:
            print(f"🎨 Generating presentation with manifest using template: {template_path}")
            
            # Load the template as base
            prs = Presentation(template_path)
            
            print(f"   Template loaded: {len(prs.slides)} slides, {len(prs.slide_layouts)} layouts")
            print(f"   Manifest provides: {len(manifest.get('layouts', []))} layout rules")
            
            # Clear all existing slides
            original_slides_count = len(prs.slides)
            slide_idxs = list(range(original_slides_count - 1, -1, -1))
            for i in slide_idxs:
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
            
            # Generate new slides using manifest rules
            slides_generated = 0
            for slide_idx, slide_content in enumerate(slide_structure.get('slides', [])):
                self._create_slide_with_manifest(prs, slide_content, manifest, slide_idx)
                slides_generated += 1
            
            print(f"   Generated {slides_generated} slides using manifest")
            
            return prs
            
        except Exception as e:
            print(f"❌ Error in manifest-based generation: {str(e)}")
            import traceback
            traceback.print_exc()
            raise e

    def generate_presentation(self, slide_structure, template_analysis, template_path):
        """Generate presentation with better template preservation (legacy method)"""
        try:
            print(f"🎨 Generating presentation using template: {template_path}")
            
            # Load the template as base - DON'T clear slides immediately
            prs = Presentation(template_path)
            
            print(f"   Template loaded: {len(prs.slides)} slides, {len(prs.slide_layouts)} layouts")
            
            # Only clear content slides, keep at least one slide as reference
            original_slides_count = len(prs.slides)
            
            # Remove slides but keep the first one as a style reference
            if original_slides_count > 1:
                slide_idxs = list(range(original_slides_count - 1, 0, -1))  # Keep first slide
                for i in slide_idxs:
                    rId = prs.slides._sldIdLst[i].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[i]
            
            # Clear content of the remaining slide but keep its layout
            if len(prs.slides) > 0:
                first_slide = prs.slides[0]
                # Clear text content but keep structure
                for shape in first_slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        shape.text = ""
            
            # Generate new slides based on structure
            slides_generated = 0
            for slide_content in slide_structure.get('slides', []):
                self._create_slide_enhanced(prs, slide_content, template_analysis)
                slides_generated += 1
            
            print(f"   Generated {slides_generated} slides")
            
            # Remove the original first slide if we generated new ones
            if slides_generated > 0 and len(prs.slides) > slides_generated:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
            
            # Save the new presentation
            output_path = os.path.join(self.temp_dir, 'generated_presentation.pptx')
            prs.save(output_path)
            
            print(f"✅ Presentation saved: {output_path}")
            
            return output_path
            
        except Exception as e:
            print(f"❌ Error generating presentation: {str(e)}")
            import traceback
            traceback.print_exc()
            raise e
    
    def _create_slide_enhanced(self, prs, slide_content, template_analysis):
        """Create slide with enhanced template styling"""
        try:
            # Choose the best layout
            layout_idx = self._choose_layout_enhanced(slide_content, template_analysis)
            layout = prs.slide_layouts[layout_idx] if layout_idx < len(prs.slide_layouts) else prs.slide_layouts[0]
            
            print(f"   Creating slide with layout: {layout.name}")
            
            slide = prs.slides.add_slide(layout)
            
            # Add title with enhanced formatting
            if 'title' in slide_content and slide.shapes.title:
                slide.shapes.title.text = slide_content['title']
                self._apply_enhanced_formatting(slide.shapes.title, template_analysis, is_title=True)
            
            # Add content with enhanced formatting
            content_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape != slide.shapes.title]
            
            if 'content' in slide_content and content_shapes:
                content_shape = content_shapes[0]
                
                if isinstance(slide_content['content'], list):
                    # Bullet points
                    content_shape.text = ""
                    for i, point in enumerate(slide_content['content']):
                        p = content_shape.text_frame.add_paragraph() if i > 0 else content_shape.text_frame.paragraphs[0]
                        p.text = point
                        p.level = 0
                else:
                    # Regular text
                    content_shape.text = slide_content['content']
                
                self._apply_enhanced_formatting(content_shape, template_analysis, is_title=False)
            
            # Add template images with better positioning
            self._add_template_images_enhanced(slide, template_analysis)
            
        except Exception as e:
            print(f"❌ Error creating slide: {str(e)}")
            import traceback
            traceback.print_exc()
    
    def _choose_layout_enhanced(self, slide_content, template_analysis):
        """Enhanced layout selection"""
        layouts = template_analysis.get('layouts', [])
        
        if not layouts:
            return 0
        
        # More sophisticated layout selection
        content_type = slide_content.get('type', '')
        has_list_content = isinstance(slide_content.get('content'), list)
        
        # Look for title slide layout for first slide
        if content_type == 'title_slide':
            for i, layout in enumerate(layouts):
                if 'title' in layout['name'].lower() and 'content' not in layout['name'].lower():
                    return i
        
        # Look for content layouts for regular slides
        if has_list_content:
            for i, layout in enumerate(layouts):
                if any('CONTENT' in p.get('type', '') for p in layout.get('placeholders', [])):
                    return i
        
        # Default to the second layout (usually title + content)
        return min(1, len(layouts) - 1)
    
    def _apply_enhanced_formatting(self, shape, template_analysis, is_title=False):
        """Apply enhanced formatting using template analysis"""
        if not shape.has_text_frame:
            return
        
        fonts = template_analysis.get('fonts', [])
        colors = template_analysis.get('colors', []) + template_analysis.get('theme_colors', [])
        
        # Apply formatting to all text
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                # Apply font with better selection
                if fonts:
                    # Use largest font for titles, medium for content
                    font_info = max(fonts, key=lambda f: f.get('size', 12)) if is_title else fonts[0]
                    
                    run.font.name = font_info.get('name', 'Calibri')
                    
                    if is_title:
                        run.font.size = Pt(max(font_info.get('size', 24), 24))
                        run.font.bold = True
                    else:
                        run.font.size = Pt(font_info.get('size', 18))
                        run.font.bold = font_info.get('bold', False)
                
                # Apply colors with better selection
                if colors:
                    # Use different colors for title vs content
                    color_idx = min(1, len(colors) - 1) if is_title else 0
                    color_hex = colors[color_idx].replace('#', '')
                    
                    if len(color_hex) == 6:
                        try:
                            rgb = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                            run.font.color.rgb = RGBColor(*rgb)
                        except:
                            pass
    
    def _add_template_images_enhanced(self, slide, template_analysis):
        """Enhanced template image addition"""
        images = template_analysis.get('images', [])
        
        if not images:
            print("   No template images found to add")
            return
        
        # Add images more intelligently
        for i, image_info in enumerate(images[:2]):  # Limit to 2 images max
            if len(slide.shapes) >= 8:  # Don't overcrowd
                break
            
            try:
                print(f"   Adding template image {i+1}")
                image_stream = io.BytesIO(image_info['image_data'])
                
                # Position images better - avoid overlapping with text
                left = image_info['left']
                top = image_info['top']
                width = image_info['width']
                height = image_info['height']
                
                # Adjust position if it conflicts with text areas
                if top < Inches(2):  # If too high, move down
                    top = Inches(3)
                
                slide.shapes.add_picture(
                    image_stream,
                    left, top, width, height
                )
                print(f"   ✅ Image added successfully")
                
            except Exception as e:
                print(f"   ❌ Failed to add image {i+1}: {str(e)}")
    
    def _get_default_analysis(self):
        """Enhanced default analysis"""
        return {
            'layouts': [],
            'colors': ['#1f4e79', '#ffffff', '#000000', '#4472c4', '#70ad47'],
            'theme_colors': ['#1f4e79', '#4472c4'],
            'fonts': [
                {'name': 'Calibri', 'size': 18, 'bold': False, 'italic': False},
                {'name': 'Calibri', 'size': 24, 'bold': True, 'italic': False}
            ],
            'images': [],
            'backgrounds': [],
            'slide_dimensions': {'width': Inches(10), 'height': Inches(7.5)},
            'master_slides': []
        }
    
    def extract_raw_template_data(self, template_path):
        """Extract raw template metadata for LLM processing"""
        try:
            prs = Presentation(template_path)
            
            raw_data = {
                "slide_size": {
                    "width_emu": int(prs.slide_width),
                    "height_emu": int(prs.slide_height)
                },
                "theme": {
                    "colors": {},
                    "fonts": {}
                },
                "layouts": [],
                "images": []
            }
            
            print(f"🔍 Extracting raw data from template: {len(prs.slide_layouts)} layouts")
            
            # Extract theme colors
            try:
                for i, master in enumerate(prs.slide_masters):
                    if master.theme and master.theme.color_scheme:
                        color_scheme = master.theme.color_scheme
                        for j in range(12):
                            try:
                                color = color_scheme[j]
                                if hasattr(color, 'rgb'):
                                    color_name = f"theme_color_{j}"
                                    if j == 0: color_name = "bg1"
                                    elif j == 1: color_name = "text1" 
                                    elif j == 2: color_name = "bg2"
                                    elif j == 3: color_name = "text2"
                                    elif j == 4: color_name = "accent1"
                                    elif j == 5: color_name = "accent2"
                                    elif j == 6: color_name = "accent3"
                                    elif j == 7: color_name = "accent4"
                                    elif j == 8: color_name = "accent5"
                                    elif j == 9: color_name = "accent6"
                                    
                                    raw_data["theme"]["colors"][color_name] = f"#{color.rgb:06x}"
                            except:
                                pass
                    
                    # Extract font scheme
                    try:
                        if master.theme and master.theme.font_scheme:
                            font_scheme = master.theme.font_scheme
                            if hasattr(font_scheme, 'major_font') and font_scheme.major_font:
                                raw_data["theme"]["fonts"]["major"] = font_scheme.major_font.latin
                            if hasattr(font_scheme, 'minor_font') and font_scheme.minor_font:
                                raw_data["theme"]["fonts"]["minor"] = font_scheme.minor_font.latin
                    except:
                        pass
            except:
                pass
            
            # Extract layouts with placeholders
            for i, layout in enumerate(prs.slide_layouts):
                layout_data = {
                    "index": i,
                    "name": layout.name,
                    "placeholders": []
                }
                
                for placeholder in layout.placeholders:
                    try:
                        ph_data = {
                            "kind": str(placeholder.placeholder_format.type).split('.')[-1],  # Get enum name
                            "left": int(placeholder.left),
                            "top": int(placeholder.top), 
                            "width": int(placeholder.width),
                            "height": int(placeholder.height)
                        }
                        layout_data["placeholders"].append(ph_data)
                    except:
                        pass
                
                raw_data["layouts"].append(layout_data)
            
            # Extract images from slides
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            img_data = {
                                "id": f"img_{slide_idx}_{shape_idx}",
                                "left": int(shape.left),
                                "top": int(shape.top),
                                "width": int(shape.width), 
                                "height": int(shape.height),
                                "usage_hint": f"appears on slide {slide_idx + 1}",
                                "size_bytes": len(shape.image.blob)
                            }
                            raw_data["images"].append(img_data)
                        except Exception as e:
                            print(f"   Could not extract image {shape_idx} from slide {slide_idx}: {e}")
            
            print(f"✅ Raw extraction complete:")
            print(f"   - Slide size: {raw_data['slide_size']}")
            print(f"   - Theme colors: {len(raw_data['theme']['colors'])}")
            print(f"   - Layouts: {len(raw_data['layouts'])}")
            print(f"   - Images: {len(raw_data['images'])}")
            
            return raw_data
            
        except Exception as e:
            print(f"❌ Error extracting raw template data: {str(e)}")
            import traceback
            traceback.print_exc()
            return {
                "slide_size": {"width_emu": 9144000, "height_emu": 6858000},
                "theme": {"colors": {}, "fonts": {}},
                "layouts": [],
                "images": []
            }

    def _create_slide_with_manifest(self, prs, slide_content, manifest, slide_idx):
        """Create slide using manifest-defined rules"""
        try:
            print(f"   Creating slide {slide_idx + 1}: {slide_content.get('title', 'Untitled')}")
            
            # Determine layout based on content type and manifest
            layout_choice = self._resolve_layout_from_manifest(slide_content, manifest, slide_idx, prs)
            
            if not layout_choice:
                print(f"     Warning: No suitable layout found, using first available")
                layout_choice = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            
            # Create slide with chosen layout
            slide = prs.slides.add_slide(layout_choice)
            print(f"     Using layout: {layout_choice.name}")
            
            # Apply manifest-based styling
            self._apply_manifest_styling(slide, slide_content, manifest)
            
            # Add template assets if specified
            self._add_manifest_assets(slide, manifest, slide_idx)
            
            # Add speaker notes if available
            self._add_speaker_notes(slide, slide_content)
            
            return slide
            
        except Exception as e:
            print(f"     Error creating slide {slide_idx + 1}: {str(e)}")
            # Fallback to basic slide creation
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            
            # Basic content population
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                slide.shapes.title.text = slide_content.get('title', '')
            
            return slide
    
    def _resolve_layout_from_manifest(self, slide_content, manifest, slide_idx, prs):
        """Choose appropriate layout based on manifest rules and content"""
        try:
            # Get layout rules from manifest
            manifest_layouts = manifest.get('layouts', [])
            
            if not manifest_layouts:
                return None
            
            # Determine content characteristics
            has_title = bool(slide_content.get('title', '').strip())
            has_content = bool(slide_content.get('content', '').strip())
            content_length = len(slide_content.get('content', ''))
            content_type = slide_content.get('layout_hint', '')
            
            # Layout selection logic - prioritize content-capable layouts
            print(f"         Layout selection: slide_idx={slide_idx}, has_content={has_content}, content_length={content_length}")
            
            # If we have content, prioritize layouts that can display content
            if has_content and content_length > 0:
                print(f"         Has content ({content_length} chars), looking for content-capable layouts...")
                
                # First priority: title_content layouts
                for layout_def in manifest_layouts:
                    if layout_def.get('archetype') == 'title_content':
                        print(f"         Found title_content layout: {layout_def.get('name')}")
                        return self._find_layout_by_name(layout_def.get('name'), prs)
                
                # Second priority: two_content layouts for longer content
                if content_length > 200:
                    for layout_def in manifest_layouts:
                        if layout_def.get('archetype') in ['two_content']:
                            print(f"         Found two_content layout: {layout_def.get('name')}")
                            return self._find_layout_by_name(layout_def.get('name'), prs)
                
                # Third priority: any content layout
                for layout_def in manifest_layouts:
                    if 'content' in layout_def.get('archetype', '').lower():
                        print(f"         Found content layout: {layout_def.get('name')}")
                        return self._find_layout_by_name(layout_def.get('name'), prs)
            
            # Only use section headers for slides explicitly marked as sections or without content
            if content_type == 'section' or not has_content:
                print(f"         Looking for section/title-only layouts...")
                for layout_def in manifest_layouts:
                    if layout_def.get('archetype') in ['section_header', 'title_only']:
                        print(f"         Found section layout: {layout_def.get('name')}")
                        return self._find_layout_by_name(layout_def.get('name'), prs)
            
            # Fallback to first layout
            if manifest_layouts:
                return self._find_layout_by_name(manifest_layouts[0].get('name'), prs)
            
            return None
            
        except Exception as e:
            print(f"     Warning: Layout resolution failed: {str(e)}")
            return None
    
    def _find_layout_by_name(self, layout_name, prs):
        """Find PowerPoint layout by name"""
        try:
            if not hasattr(prs, 'slide_layouts'):
                return None
                
            for layout in prs.slide_layouts:
                if layout.name == layout_name:
                    return layout
            
            # Fallback - return first available layout
            return prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            
        except:
            return None
    
    def _apply_manifest_styling(self, slide, slide_content, manifest):
        """Apply text formatting and colors from manifest"""
        try:
            # Get styling rules from manifest
            text_defaults = manifest.get('text_defaults', {})
            theme_palette = manifest.get('theme', {}).get('palette', {})
            
            print(f"       Debug: slide_content = {slide_content}")
            
            # Apply title styling
            title_text = slide_content.get('title', '')
            print(f"       Title: '{title_text}'")
            
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title_shape = slide.shapes.title
                title_shape.text = title_text
                
                if title_shape.text_frame:
                    title_rules = text_defaults.get('title', {})
                    self._apply_text_formatting(title_shape.text_frame, title_rules, theme_palette)
            
            # Apply body content styling
            body_content = slide_content.get('content', '')
            print(f"       Content: '{body_content}' (type: {type(body_content)})")
            
            # Handle empty content with fallback
            if not body_content:
                print(f"       ⚠️  Empty content detected, using title as content fallback")
                body_content = f"Key points about {title_text}" if title_text else "Content will be added here"
            
            if body_content:
                # Handle list content vs string content
                if isinstance(body_content, list):
                    body_text = '\n'.join(f"• {item}" for item in body_content)
                else:
                    body_text = str(body_content)
                
                print(f"       Processed body text: '{body_text[:100]}...'")
                
                # Find body placeholder - try multiple approaches
                body_shape = None
                
                # Method 1: Find by placeholder type (BODY = 2)
                for shape in slide.shapes:
                    if hasattr(shape, 'placeholder_format'):
                        ph_type = shape.placeholder_format.type
                        print(f"         Found placeholder type: {ph_type}")
                        if ph_type == 2:  # BODY placeholder
                            body_shape = shape
                            print(f"         Using BODY placeholder (type 2)")
                            break
                
                # Method 2: Find by placeholder index if method 1 fails
                if not body_shape:
                    try:
                        print(f"         Trying placeholder by index...")
                        for i, shape in enumerate(slide.placeholders):
                            print(f"           Placeholder {i}: {shape.placeholder_format.type if hasattr(shape, 'placeholder_format') else 'no format'}")
                            if i == 1:  # Usually second placeholder is body
                                body_shape = shape
                                print(f"         Using placeholder index 1")
                                break
                    except Exception as e:
                        print(f"         Placeholder index method failed: {e}")
                
                # Method 3: Find text placeholders that aren't title
                if not body_shape:
                    print(f"         Looking for any non-title text placeholder...")
                    for shape in slide.shapes:
                        if hasattr(shape, 'placeholder_format') and hasattr(shape, 'text_frame'):
                            ph_type = shape.placeholder_format.type
                            if ph_type != 1 and shape != getattr(slide.shapes, 'title', None):  # Not TITLE placeholder
                                body_shape = shape
                                print(f"         Using non-title placeholder (type {ph_type})")
                                break
                
                # Method 4: Find any text shape that's not the title
                if not body_shape:
                    print(f"         Using fallback: any text shape that's not title...")
                    for shape in slide.shapes:
                        if (hasattr(shape, 'text_frame') and 
                            shape != getattr(slide.shapes, 'title', None)):
                            body_shape = shape
                            print(f"         Using fallback text shape")
                            break
                
                if body_shape and hasattr(body_shape, 'text_frame'):
                    body_shape.text = body_text
                    body_rules = text_defaults.get('body', [{}])[0]
                    self._apply_text_formatting(body_shape.text_frame, body_rules, theme_palette)
                    print(f"       ✅ Content applied to body shape")
                else:
                    print(f"       ❌ No suitable body shape found")
                    print(f"         Available shapes: {[type(s).__name__ for s in slide.shapes]}")
            else:
                print(f"       ⚠️ No body content to apply")
            
        except Exception as e:
            print(f"     ❌ Styling application failed: {str(e)}")
            import traceback
            traceback.print_exc()
    
    def _apply_text_formatting(self, text_frame, format_rules, palette):
        """Apply specific formatting rules to text frame"""
        try:
            if not text_frame.paragraphs:
                return
            
            paragraph = text_frame.paragraphs[0]
            if not paragraph.runs:
                return
            
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            font = run.font
            
            # Apply font family
            if format_rules.get('family'):
                font.name = format_rules['family']
            
            # Apply font size
            if format_rules.get('size_pt'):
                font.size = Pt(format_rules['size_pt'])
            
            # Apply bold
            if format_rules.get('bold'):
                font.bold = True
            
            # Apply color from palette
            color_key = format_rules.get('color')
            if color_key and palette.get(color_key):
                color_hex = palette[color_key].replace('#', '')
                if len(color_hex) == 6:
                    try:
                        rgb_color = RGBColor(
                            int(color_hex[0:2], 16),
                            int(color_hex[2:4], 16), 
                            int(color_hex[4:6], 16)
                        )
                        font.color.rgb = rgb_color
                    except:
                        pass
            
        except Exception as e:
            print(f"     Warning: Text formatting failed: {str(e)}")
    
    def _add_manifest_assets(self, slide, manifest, slide_idx):
        """Add template images/assets based on manifest rules"""
        try:
            assets = manifest.get('assets', [])
            
            for asset in assets:
                apply_policy = asset.get('apply_on', 'none')
                
                # Check if this asset should be applied to this slide
                should_apply = (
                    apply_policy == 'all' or
                    (apply_policy == 'title_only' and slide_idx == 0)
                )
                
                if should_apply:
                    # For now, just print that we would add this asset
                    print(f"     Would add asset {asset.get('id')} at position ({asset.get('left')}, {asset.get('top')})")
                    # TODO: Implement actual image copying from template
            
        except Exception as e:
            print(f"     Warning: Asset addition failed: {str(e)}")

    def _add_speaker_notes(self, slide, slide_content):
        """Add speaker notes to the slide"""
        try:
            speaker_notes = slide_content.get('speaker_notes', '')
            
            if speaker_notes:
                print(f"       Adding speaker notes: {len(speaker_notes)} characters")
                
                # Access the notes slide for this slide
                notes_slide = slide.notes_slide
                if notes_slide and hasattr(notes_slide, 'notes_text_frame'):
                    notes_slide.notes_text_frame.text = speaker_notes
                    print(f"       ✅ Speaker notes added successfully")
                else:
                    print(f"       ⚠️  Notes slide not accessible")
            else:
                print(f"       No speaker notes provided for this slide")
                
        except Exception as e:
            print(f"     Warning: Speaker notes addition failed: {str(e)}")

    def cleanup(self):
        """Clean up temporary files"""
        try:
            import shutil
            shutil.rmtree(self.temp_dir)
        except:
            pass

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import tempfile
import json
from PIL import Image
import io

class PPTProcessor:
    def __init__(self):
        self.temp_dir = tempfile.mkdtemp()
    
    def analyze_template(self, template_path):
        """
        Analyze the uploaded PowerPoint template to extract:
        - Slide layouts
        - Color schemes
        - Font styles
        - Images and their positions
        - Master slide properties
        """
        try:
            prs = Presentation(template_path)
            analysis = {
                'layouts': [],
                'colors': [],
                'fonts': [],
                'images': [],
                'slide_dimensions': {
                    'width': prs.slide_width,
                    'height': prs.slide_height
                },
                'master_slides': []
            }
            
            # Analyze slide layouts
            for i, layout in enumerate(prs.slide_layouts):
                layout_info = {
                    'index': i,
                    'name': layout.name,
                    'placeholders': []
                }
                
                # Analyze placeholders in each layout
                for placeholder in layout.placeholders:
                    placeholder_info = {
                        'index': placeholder.placeholder_format.idx,
                        'type': str(placeholder.placeholder_format.type),
                        'name': placeholder.name,
                        'left': placeholder.left,
                        'top': placeholder.top,
                        'width': placeholder.width,
                        'height': placeholder.height
                    }
                    layout_info['placeholders'].append(placeholder_info)
                
                analysis['layouts'].append(layout_info)
            
            # Analyze existing slides for styling
            for slide in prs.slides:
                self._extract_slide_styles(slide, analysis)
            
            # Extract theme colors from master slides
            for master in prs.slide_masters:
                master_info = self._extract_master_styles(master)
                analysis['master_slides'].append(master_info)
            
            return analysis
            
        except Exception as e:
            print(f"Error analyzing template: {str(e)}")
            return self._get_default_analysis()
    
    def _extract_slide_styles(self, slide, analysis):
        """Extract styling information from individual slides"""
        for shape in slide.shapes:
            # Extract text formatting
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name and run.font.name not in analysis['fonts']:
                            analysis['fonts'].append({
                                'name': run.font.name,
                                'size': run.font.size.pt if run.font.size else 18,
                                'bold': run.font.bold,
                                'italic': run.font.italic
                            })
                        
                        # Extract colors
                        if run.font.color.rgb:
                            color = f"#{run.font.color.rgb:06x}"
                            if color not in analysis['colors']:
                                analysis['colors'].append(color)
            
            # Extract images
            if shape.shape_type == 13:  # Picture shape type
                try:
                    image_data = {
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'image_data': shape.image.blob
                    }
                    analysis['images'].append(image_data)
                except:
                    pass
    
    def _extract_master_styles(self, master):
        """Extract styling from master slides"""
        master_info = {
            'name': master.name,
            'colors': [],
            'fonts': []
        }
        
        # Extract theme colors
        try:
            color_scheme = master.theme.color_scheme
            for i in range(12):  # Standard Office color scheme has 12 colors
                try:
                    color = color_scheme[i]
                    master_info['colors'].append(f"#{color.rgb:06x}")
                except:
                    pass
        except:
            pass
        
        return master_info
    
    def _get_default_analysis(self):
        """Fallback analysis if template analysis fails"""
        return {
            'layouts': [],
            'colors': ['#1f4e79', '#ffffff', '#000000', '#4472c4'],
            'fonts': [{'name': 'Calibri', 'size': 18, 'bold': False, 'italic': False}],
            'images': [],
            'slide_dimensions': {'width': Inches(10), 'height': Inches(7.5)},
            'master_slides': []
        }
    
    def generate_presentation(self, slide_structure, template_analysis, template_path):
        """
        Generate a new PowerPoint presentation based on:
        - slide_structure: LLM-generated content structure
        - template_analysis: Extracted template styling
        - template_path: Original template file for layout copying
        """
        try:
            # Load the template as base
            prs = Presentation(template_path)
            
            # Clear existing slides but keep master
            slide_idxs = list(range(len(prs.slides) - 1, -1, -1))
            for i in slide_idxs:
                rId = prs.slides._sldIdLst[i].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[i]
            
            # Generate slides based on structure
            for slide_content in slide_structure.get('slides', []):
                self._create_slide(prs, slide_content, template_analysis)
            
            # Save the new presentation
            output_path = os.path.join(self.temp_dir, 'generated_presentation.pptx')
            prs.save(output_path)
            
            return output_path
            
        except Exception as e:
            print(f"Error generating presentation: {str(e)}")
            raise e
    
    def _create_slide(self, prs, slide_content, template_analysis):
        """Create a single slide with the given content"""
        try:
            # Choose appropriate layout based on content type
            layout_idx = self._choose_layout(slide_content, template_analysis)
            layout = prs.slide_layouts[layout_idx] if layout_idx < len(prs.slide_layouts) else prs.slide_layouts[0]
            
            slide = prs.slides.add_slide(layout)
            
            # Add title
            if 'title' in slide_content and slide.shapes.title:
                slide.shapes.title.text = slide_content['title']
                self._apply_text_formatting(slide.shapes.title, template_analysis, is_title=True)
            
            # Add content
            content_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape != slide.shapes.title]
            
            if 'content' in slide_content and content_shapes:
                content_shape = content_shapes[0]
                
                if isinstance(slide_content['content'], list):
                    # Bullet points
                    content_shape.text = ""
                    for i, point in enumerate(slide_content['content']):
                        p = content_shape.text_frame.add_paragraph() if i > 0 else content_shape.text_frame.paragraphs[0]
                        p.text = point
                        p.level = 0
                else:
                    # Regular text
                    content_shape.text = slide_content['content']
                
                self._apply_text_formatting(content_shape, template_analysis, is_title=False)
            
            # Add images from template if available
            self._add_template_images(slide, template_analysis)
            
        except Exception as e:
            print(f"Error creating slide: {str(e)}")
    
    def _choose_layout(self, slide_content, template_analysis):
        """Choose the most appropriate layout for the content"""
        layouts = template_analysis.get('layouts', [])
        
        if not layouts:
            return 0  # Default to first layout
        
        # Simple heuristic: choose based on content type
        if 'content' in slide_content and isinstance(slide_content['content'], list):
            # Bullet points - prefer layouts with content placeholders
            for i, layout in enumerate(layouts):
                if any('CONTENT' in p.get('type', '') for p in layout.get('placeholders', [])):
                    return i
        
        # Default to title and content layout (usually index 1)
        return min(1, len(layouts) - 1)
    
    def _apply_text_formatting(self, shape, template_analysis, is_title=False):
        """Apply template-based formatting to text"""
        if not shape.has_text_frame:
            return
        
        fonts = template_analysis.get('fonts', [])
        colors = template_analysis.get('colors', [])
        
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                # Apply font
                if fonts:
                    font_info = fonts[0]  # Use first font as default
                    run.font.name = font_info.get('name', 'Calibri')
                    
                    if is_title:
                        run.font.size = Pt(font_info.get('size', 28) + 10)
                        run.font.bold = True
                    else:
                        run.font.size = Pt(font_info.get('size', 18))
                        run.font.bold = font_info.get('bold', False)
                
                # Apply color
                if colors and len(colors) > 1:
                    color_hex = colors[1 if is_title else 0].replace('#', '')
                    if len(color_hex) == 6:
                        try:
                            rgb = tuple(int(color_hex[i:i+2], 16) for i in (0, 2, 4))
                            run.font.color.rgb = RGBColor(*rgb)
                        except:
                            pass
    
    def _add_template_images(self, slide, template_analysis):
        """Add images from the template to appropriate positions"""
        images = template_analysis.get('images', [])
        
        # Add first template image if available and there's space
        if images and len(slide.shapes) < 5:  # Avoid overcrowding
            try:
                image_info = images[0]
                image_stream = io.BytesIO(image_info['image_data'])
                
                slide.shapes.add_picture(
                    image_stream,
                    image_info['left'],
                    image_info['top'],
                    image_info['width'],
                    image_info['height']
                )
            except Exception as e:
                print(f"Error adding template image: {str(e)}")
    
    def cleanup(self):
        """Clean up temporary files"""
        try:
            import shutil
            shutil.rmtree(self.temp_dir)
        except:
            pass

